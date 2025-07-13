"""
Simple Document Processor - Direct Text Extraction & Processing
No LlamaIndex dependencies - pure Python text extraction for experimentation
"""

import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional
from uuid import UUID
from dataclasses import dataclass

from sqlalchemy.ext.asyncio import AsyncSession

from src.backend.models.database import File
from src.backend.config.settings import get_settings

settings = get_settings()


@dataclass
class ProcessedDocument:
    """Result of document processing"""
    file_id: UUID
    chunks_created: int
    processing_method: str
    success: bool
    error_message: Optional[str] = None


class SimpleDocumentProcessor:
    """
    Direct document processor for embedding experimentation
    
    Features:
    - Pure Python text extraction
    - Multiple file format support
    - Direct pgvector storage
    - Perfect for chunking experiments
    """
    
    def __init__(self, db: AsyncSession, rag_service):
        self.db = db
        self.rag_service = rag_service
    
    async def process_file(self, file_record: File, file_path: Path) -> ProcessedDocument:
        """Process file with direct text extraction"""
        try:
            # Extract text based on file type
            text_content = await self._extract_text(file_path)
            
            if not text_content or len(text_content.strip()) < 10:
                return ProcessedDocument(
                    file_id=file_record.id,
                    chunks_created=0,
                    processing_method="text_extraction_failed",
                    success=False,
                    error_message="No meaningful text content extracted"
                )
            
            # Add to RAG service (handles chunking and embedding)
            success = await self.rag_service.add_document(
                file_content=text_content,
                tenant_slug=file_record.tenant_slug,
                file_record=file_record
            )
            
            if success:
                # Count chunks created (simple estimate)
                word_count = len(text_content.split())
                estimated_chunks = max(1, word_count // 400)  # Rough estimate
                
                return ProcessedDocument(
                    file_id=file_record.id,
                    chunks_created=estimated_chunks,
                    processing_method="direct_text_extraction",
                    success=True
                )
            else:
                return ProcessedDocument(
                    file_id=file_record.id,
                    chunks_created=0,
                    processing_method="embedding_failed",
                    success=False,
                    error_message="Failed to create embeddings"
                )
                
        except Exception as e:
            return ProcessedDocument(
                file_id=file_record.id,
                chunks_created=0,
                processing_method="error",
                success=False,
                error_message=str(e)
            )
    
    async def _extract_text(self, file_path: Path) -> str:
        """Extract text from various file formats"""
        file_extension = file_path.suffix.lower()
        
        try:
            if file_extension == '.txt':
                return await self._extract_txt(file_path)
            elif file_extension == '.pdf':
                return await self._extract_pdf(file_path)
            elif file_extension in ['.doc', '.docx']:
                return await self._extract_docx(file_path)
            elif file_extension in ['.xls', '.xlsx']:
                return await self._extract_excel(file_path)
            elif file_extension in ['.ppt', '.pptx']:
                return await self._extract_pptx(file_path)
            elif file_extension in ['.html', '.htm']:
                return await self._extract_html(file_path)
            else:
                # Try as plain text
                return await self._extract_txt(file_path)
                
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""
    
    async def _extract_txt(self, file_path: Path) -> str:
        """Extract from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except:
                    continue
            return ""
    
    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract from PDF files using PyPDF2"""
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except ImportError:
            print("PyPDF2 not available for PDF processing")
            return ""
        except Exception as e:
            print(f"Error extracting PDF: {e}")
            return ""
    
    async def _extract_docx(self, file_path: Path) -> str:
        """Extract from Word documents using python-docx"""
        try:
            import docx
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            print("python-docx not available for DOCX processing")
            return ""
        except Exception as e:
            print(f"Error extracting DOCX: {e}")
            return ""
    
    async def _extract_excel(self, file_path: Path) -> str:
        """Extract from Excel files using openpyxl"""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text
        except ImportError:
            print("openpyxl not available for Excel processing")
            return ""
        except Exception as e:
            print(f"Error extracting Excel: {e}")
            return ""
    
    async def _extract_pptx(self, file_path: Path) -> str:
        """Extract from PowerPoint files using python-pptx"""
        try:
            import pptx
            presentation = pptx.Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(presentation.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except ImportError:
            print("python-pptx not available for PPTX processing")
            return ""
        except Exception as e:
            print(f"Error extracting PPTX: {e}")
            return ""
    
    async def _extract_html(self, file_path: Path) -> str:
        """Extract from HTML files"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text()
        except ImportError:
            print("BeautifulSoup not available for HTML processing")
            return ""
        except Exception as e:
            print(f"Error extracting HTML: {e}")
            return ""